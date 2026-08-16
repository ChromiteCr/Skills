#!/usr/bin/env python3
"""从片单 JSON 生成可编辑的 .pptx。

    python3 outline_to_pptx.py deck.json [out.pptx]

这条路的价值是「能在 Keynote 或 PowerPoint 里继续改」。**它会丢东西，要如实告诉使用者**：

  - python-pptx 完全不支持动画与转场，导出来的片子没有任何动效
  - 底色只能是纯色，做不到 HTML 模板里那种多层渐变
  - 颗粒、圆角容器、bento 的大小分层都被简化

要视觉精度就用 templates/deck.html；要可编辑就用这个。两者不冲突，可以都给。

字号下限 30pt，与 HTML 模板里的 40px 是同一条规则在两种单位下的写法
（30pt = 40px，96dpi）。塞不下就删字，不要调低。

片单格式（type 决定版式，其余字段按类型取用）：

    {
      "accent": "#5a8dee",
      "theme": "dark",
      "font": "Helvetica Neue",
      "slides": [
        {"type": "title",   "value": "产品名", "caption": "一句话定位"},
        {"type": "phrase",  "value": "转账不该点七次"},
        {"type": "num",     "value": "7 次", "caption": "完成一笔转账",
                            "unverified": false},
        {"type": "section", "value": "章节名"},
        {"type": "feature", "value": "功能名", "caption": "一句话"},
        {"type": "specs",   "value": "规格",
                            "items": [["项", "值"], ["项", "值"]]},
        {"type": "versus",  "value": "对比什么",
                            "items": [["上一代", "1.0"], ["本代", "3.2"]]},
        {"type": "price",   "value": "产品名",
                            "items": [["¥4999", "128GB"]], "caption": "9 月 20 日"},
        {"type": "close",   "value": "收束一句"}
      ]
    }
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt
except ImportError:
    sys.exit("缺 python-pptx。装：pip3 install python-pptx")

# 16:9
W, H = Inches(13.333), Inches(7.5)
# 内容离边 10–15%
PAD_X, PAD_Y = Inches(1.6), Inches(0.9)
BODY_W = W - PAD_X * 2

# 字号：下限 30pt，不可破
SIZE = {"num": 160, "phrase": 90, "feature": 60, "title": 48,
        "sub": 30, "label": 30}

THEME = {
    "dark":  {"bg": "0A0B0E", "ink": "F5F5F7", "muted": "8E8E93"},
    "light": {"bg": "F5F5F7", "ink": "1D1D1F", "muted": "6E6E73"},
}


def rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#").upper())


class Deck:
    def __init__(self, spec: dict):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]
        pal = THEME[spec.get("theme", "dark")]
        self.bg = rgb(pal["bg"])
        self.ink = rgb(pal["ink"])
        self.muted = rgb(pal["muted"])
        self.accent = rgb(spec.get("accent", "#5A8DEE"))
        self.font = spec.get("font", "Helvetica Neue")

    def slide(self):
        s = self.prs.slides.add_slide(self.blank)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg
        return s

    def text(self, slide, body, *, top, height, size, color=None,
             bold=False, align=PP_ALIGN.CENTER, left=None, width=None):
        box = slide.shapes.add_textbox(
            left if left is not None else PAD_X,
            top,
            width if width is not None else BODY_W,
            height,
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(body)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = self.font
        run.font.color.rgb = color or self.ink
        return box

    # ── 各片型 ────────────────────────────────────────────
    def title(self, d):
        s = self.slide()
        self.text(s, d["value"], top=Inches(2.4), height=Inches(1.9),
                  size=SIZE["phrase"], bold=True)
        if d.get("caption"):
            self.text(s, d["caption"], top=Inches(4.5), height=Inches(0.8),
                      size=SIZE["sub"], color=self.muted)

    def phrase(self, d):
        s = self.slide()
        self.text(s, d["value"], top=Inches(2.5), height=Inches(2.5),
                  size=SIZE["phrase"], bold=True)

    def num(self, d):
        s = self.slide()
        self.text(s, d["value"], top=Inches(2.0), height=Inches(2.6),
                  size=SIZE["num"], bold=True, color=self.accent)
        if d.get("caption"):
            self.text(s, d["caption"], top=Inches(4.7), height=Inches(0.7),
                      size=SIZE["sub"], color=self.muted)
        # 未确认项在片上必须看得见，不能只写在大纲里
        if d.get("unverified"):
            self.text(s, "未确认 · 材料里没有出处", top=Inches(5.5),
                      height=Inches(0.6), size=SIZE["label"], color=self.muted)

    def section(self, d):
        s = self.slide()
        self.text(s, d["value"], top=Inches(3.1), height=Inches(1.3),
                  size=SIZE["title"], color=self.muted)

    def feature(self, d):
        s = self.slide()
        self.text(s, d["value"], top=Inches(2.6), height=Inches(1.4),
                  size=SIZE["feature"], bold=True)
        if d.get("caption"):
            self.text(s, d["caption"], top=Inches(4.2), height=Inches(0.8),
                      size=SIZE["sub"], color=self.muted)

    def specs(self, d):
        s = self.slide()
        self.text(s, d["value"], top=PAD_Y, height=Inches(1.0),
                  size=SIZE["title"], bold=True, align=PP_ALIGN.LEFT)
        items = d.get("items", [])[:6]
        col_w = BODY_W / 3
        for n, (k, v) in enumerate(items):
            cx = PAD_X + col_w * (n % 3)
            cy = Inches(2.6) + Inches(1.9) * (n // 3)
            self.text(s, k, top=cy, height=Inches(0.55), size=SIZE["label"],
                      color=self.muted, align=PP_ALIGN.LEFT,
                      left=cx, width=col_w - Inches(0.3))
            self.text(s, v, top=cy + Inches(0.6), height=Inches(0.9), size=64,
                      bold=True, align=PP_ALIGN.LEFT,
                      left=cx, width=col_w - Inches(0.3))

    def versus(self, d):
        s = self.slide()
        self.text(s, d["value"], top=PAD_Y, height=Inches(1.0),
                  size=SIZE["title"], bold=True, align=PP_ALIGN.LEFT)
        items = d.get("items", [])[:2]
        col_w = BODY_W / 2
        for n, (head, val) in enumerate(items):
            cx = PAD_X + col_w * n
            # 列头必须标明对比的轴，否则「快 3 倍」是没有对象的说法
            self.text(s, head, top=Inches(2.6), height=Inches(0.7),
                      size=SIZE["sub"], color=self.muted,
                      left=cx, width=col_w - Inches(0.4))
            self.text(s, val, top=Inches(3.4), height=Inches(1.8), size=100,
                      bold=True, color=self.accent if n else self.ink,
                      left=cx, width=col_w - Inches(0.4))

    def price(self, d):
        s = self.slide()
        self.text(s, d["value"], top=Inches(1.4), height=Inches(1.0),
                  size=SIZE["title"], bold=True)
        items = d.get("items", [])[:3]
        if items:
            col_w = BODY_W / len(items)
            for n, (p, cfg) in enumerate(items):
                cx = PAD_X + col_w * n
                self.text(s, p, top=Inches(3.0), height=Inches(1.4), size=90,
                          bold=True, color=self.accent if n == 0 else self.ink,
                          left=cx, width=col_w)
                self.text(s, cfg, top=Inches(4.4), height=Inches(0.7),
                          size=SIZE["sub"], color=self.muted,
                          left=cx, width=col_w)
        if d.get("caption"):
            self.text(s, d["caption"], top=Inches(5.4), height=Inches(0.7),
                      size=SIZE["sub"], color=self.muted)

    def close(self, d):
        self.phrase(d)


def main() -> int:
    if len(sys.argv) < 2:
        sys.exit(__doc__.split("片单格式")[0].strip())

    src = Path(sys.argv[1])
    if not src.is_file():
        sys.exit(f"找不到文件: {src}")

    spec = json.loads(src.read_text(encoding="utf-8"))
    slides = spec.get("slides") or []
    if not slides:
        sys.exit("片单里没有 slides")

    deck = Deck(spec)
    handlers = {name: getattr(deck, name) for name in
                ("title", "phrase", "num", "section", "feature",
                 "specs", "versus", "price", "close")}

    unknown, made = set(), 0
    for d in slides:
        fn = handlers.get(d.get("type"))
        if fn is None:
            unknown.add(d.get("type"))
            continue
        fn(d)
        made += 1

    out = Path(sys.argv[2]) if len(sys.argv) > 2 else src.with_suffix(".pptx")
    deck.prs.save(str(out))

    print(f"已生成: {out}  ({made} 页)")
    if unknown:
        print(f"跳过了不认识的片型: {', '.join(sorted(map(str, unknown)))}")
    print("这份 pptx 没有动画与转场，底色是纯色 —— python-pptx 做不到这些。")
    print("要那个层次的视觉就用 templates/deck.html。")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
