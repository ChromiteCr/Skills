#!/usr/bin/env python3
"""从现有 .pptx 里提取文字与结构，供重构成发布会风格时读。

    python3 read_pptx.py deck.pptx

只读，不改原文件。输出每页的文字、表格、图片数量，以及一个「这页太密」的提示：
一页超过 25 个词或出现项目符号列表时标出来，这些是需要拆片的地方。

为什么要有这个脚本：把整份 pptx 的 XML 塞进上下文很贵，而且大部分是版式噪音。
这里只取文字和密度信号。
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("缺 python-pptx。装：pip3 install python-pptx")

# 该拆片的判据。光数总词数会冤枉规格密排和 bento——那两种片型本来就有很多标签值对，
# 但每条都很短。真正需要拆的信号是片上出现了整句散文，或者有项目符号层级。
PROSE_WORDS = 20   # 单行超过这个词数，基本是一整句，属于该留给讲的人说的部分
TOTAL_WORDS = 45   # 总量上限，标签值对再多也不该超过这个


def words(text: str) -> int:
    """中英混排的粗略计数：非 ASCII 按字算，ASCII 按空格分词。"""
    cjk = sum(1 for ch in text if ord(ch) > 0x2E80)
    latin = len([w for w in "".join(
        ch if ord(ch) <= 0x2E80 else " " for ch in text).split() if w])
    return cjk + latin


def main() -> int:
    if len(sys.argv) != 2:
        sys.exit("用法: python3 read_pptx.py <file.pptx>")

    path = Path(sys.argv[1])
    if not path.is_file():
        sys.exit(f"找不到文件: {path}")

    prs = Presentation(str(path))
    w, h = prs.slide_width, prs.slide_height
    ratio = f"{w / h:.3f}" if h else "?"
    print(f"# {path.name}")
    print(f"尺寸 {Emu(w).inches:.2f}in × {Emu(h).inches:.2f}in · 宽高比 {ratio} "
          f"({'16:9' if abs(w / h - 16 / 9) < 0.01 else '不是 16:9，重构时要改'})")
    print(f"共 {len(prs.slides)} 页\n")

    total_dense = 0

    for n, slide in enumerate(prs.slides, 1):
        lines, tables, pics, bullets = [], 0, 0, 0

        for shape in slide.shapes:
            if shape.shape_type == 13 or shape.__class__.__name__ == "Picture":
                pics += 1
            if getattr(shape, "has_table", False) and shape.has_table:
                tables += 1
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if not text:
                    continue
                # level > 0 基本就是项目符号层级
                if para.level > 0:
                    bullets += 1
                lines.append(text)

        counts = [words(t) for t in lines]
        wc = sum(counts)
        prose = sum(1 for c in counts if c > PROSE_WORDS)

        reasons = []
        if prose:
            reasons.append(f"{prose} 行是整句")
        if bullets:
            reasons.append(f"项目符号 {bullets} 行")
        if wc > TOTAL_WORDS:
            reasons.append(f"总量 {wc} 词")
        if reasons:
            total_dense += 1

        print(f"## 第 {n} 页 · {wc} 词" +
              (f" · 图 {pics}" if pics else "") +
              (f" · 表 {tables}" if tables else "") +
              (f"  ← 需要拆：{'，'.join(reasons)}" if reasons else ""))
        for t, c in zip(lines, counts):
            print(f"  - {t}" + (f"   ← {c} 词，整句，该留给讲的人说"
                                if c > PROSE_WORDS else ""))
        print()

    print(f"---\n{total_dense} / {len(prs.slides)} 页需要拆。")
    print("整句的那几行：片上只留碎片、数字或图，句子留给讲的人说。")
    print("项目符号那几行：顺序性的拆片，并列性的改成 bento 格或图形。")
    print("规格密排与 bento 本来就有很多短标签，词数高不算问题，看的是有没有整句。")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
